import streamlit as st
import pandas as pd
import PyPDF2
import docx
import re
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.stem import WordNetLemmatizer
from collections import Counter
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
import spacy
import matplotlib.pyplot as plt
import seaborn as sns
import io
import google.generativeai as genai
import json
import os
from datetime import datetime
import extract_msg
import pptx
import numpy as np
import easyocr
import cv2
from PIL import Image
import warnings
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.opc.packuri import PackURI
import zipfile
import base64
warnings.filterwarnings('ignore', category=UserWarning)

# Set page configuration first
st.set_page_config(
    page_title="DocInsight Pro - Advanced Document Analyzer",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for professional look
st.markdown("""
<style>
    /* Main theme colors and styling */
    :root {
        --primary-color: #4F8BF9;
        --accent-color: #1E3A8A;
        --background-color: #F8F9FA;
        --text-color: #1E293B;
        --sidebar-color: #EFF6FF;
        --card-color: white;
        --success-color: #10B981;
        --warning-color: #F59E0B;
        --error-color: #EF4444;
    }
    
    /* Body styling */
    .main {
        background-color: var(--background-color);
        color: var(--text-color);
    }
    
    /* Headers */
    h1 {
        color: var(--accent-color);
        font-weight: 700;
        margin-bottom: 1rem;
    }
    
    h2, h3 {
        color: var(--accent-color);
        font-weight: 600;
    }
    
    /* Cards/Sections styling */
    .stTabs [data-baseweb="tab-panel"] {
        background-color: var(--card-color);
        border-radius: 10px;
        padding: 1rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05);
        margin-top: 1rem;
    }
    
    /* Sidebar styling */
    .css-1d391kg {
        background-color: var(--sidebar-color);
    }
    
    /* Buttons */
    .stButton>button {
        background-color: var(--primary-color);
        color: white;
        border-radius: 6px;
        padding: 0.5rem 1rem;
        border: none;
        font-weight: 500;
        transition: all 0.3s;
    }
    .stButton>button:hover {
        background-color: var(--accent-color);
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
    }
    
    /* Metrics */
    .css-50ug3q {
        background-color: var(--card-color);
        border-radius: 8px;
        padding: 1rem;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.05);
    }
    
    /* File uploader */
    .css-1cpxqw2 {
        background-color: var(--card-color);
        border-radius: 8px;
        padding: 1rem;
        border: 2px dashed var(--primary-color);
    }
    
    /* Expanders */
    .streamlit-expanderHeader {
        background-color: var(--card-color);
        border-radius: 8px;
    }
    
    /* Success/Warning/Error messages */
    .element-container .stAlert {
        border-radius: 8px;
    }

    /* Fix file uploader alignment */
    .stFileUploader > div {
        width: 100% !important;
    }
    
    /* Fix document type cards layout */
    .document-type-card {
        background-color: #EFF6FF;
        padding: 15px;
        border-radius: 10px;
        text-align: center;
        height: 100%;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
    }
    
    /* Fix cards spacing and alignment */
    .feature-card {
        background-color: white;
        padding: 20px;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
        height: 100%;
        display: flex;
        flex-direction: column;
    }
    
    /* Fix icon and text alignment */
    .icon-heading {
        display: flex;
        align-items: center;
        gap: 8px;
        margin-bottom: 10px;
    }
    
    /* Better sidebar spacing */
    .sidebar-content {
        margin: 20px 0;
    }
    
    /* Hide Streamlit's default footer */
    footer {
        visibility: hidden;
    }
    
    /* Make the main header more prominent */
    .main-header {
        font-size: 2.2rem;
        font-weight: 700;
        color: var(--accent-color);
        margin-bottom: 0.5rem;
        text-align: center;
    }
    
    /* Better subheader */
    .main-subheader {
        font-size: 1.1rem;
        color: var(--text-color);
        opacity: 0.8;
        margin-bottom: 2rem;
        text-align: center;
    }
    
    /* Image gallery styling */
    .image-gallery {
        display: flex;
        flex-wrap: wrap;
        gap: 10px;
        margin-top: 10px;
    }
    
    .image-thumbnail {
        border: 1px solid #ddd;
        border-radius: 5px;
        padding: 5px;
        cursor: pointer;
        transition: transform 0.3s;
    }
    
    .image-thumbnail:hover {
        transform: scale(1.05);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    
    /* OCR results table */
    .ocr-table {
        width: 100%;
        border-collapse: collapse;
        margin-top: 15px;
    }
    
    .ocr-table th, .ocr-table td {
        border: 1px solid #ddd;
        padding: 8px;
        text-align: left;
    }
    
    .ocr-table th {
        background-color: #f1f8ff;
        color: var(--accent-color);
    }
    
    .ocr-table tr:nth-child(even) {
        background-color: #f9f9f9;
    }
</style>
""", unsafe_allow_html=True)

# Download necessary NLTK data
@st.cache_resource
def load_ocr_reader():
    """Initialize EasyOCR reader"""
    try:
        reader = easyocr.Reader(['en'])
        return reader
    except Exception as e:
        st.error(f"Error loading EasyOCR: {e}")
        return None

def download_nltk_data():
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)
    nltk.download('wordnet', quiet=True)

download_nltk_data()

# Load spaCy model for NER
@st.cache_resource
def load_spacy_model():
    try:
        return spacy.load("en_core_web_sm")
    except:
        st.warning("Please install the spaCy model: python -m spacy download en_core_web_sm")
        return None

nlp = load_spacy_model()

def extract_text_from_pdf(file, extract_images_text=True):
    """Extract text from PDF with optional image OCR extraction"""
    try:
        # Create a bytes buffer
        file_bytes = io.BytesIO(file.getvalue())
        
        # Try to read PDF with PyPDF2
        pdf_reader = PyPDF2.PdfReader(file_bytes, strict=False)
        text = ""
        
        # Dictionary to store extracted images and their OCR results
        extracted_image_data = []
        
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                # Extract text from page
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                
                # If image extraction is enabled
                if extract_images_text:
                    # Get images from this page
                    page_images = extract_images_from_pdf_page(file, page_num)
                    
                    # If images were found on this page, process them with OCR
                    if page_images:
                        for i, img in enumerate(page_images):
                            try:
                                ocr_result = extract_text_from_image(img)
                                if isinstance(ocr_result, dict) and ocr_result['text'].strip():
                                    # Add image OCR text to the PDF text
                                    text += f"\n[IMAGE TEXT (Page {page_num+1}, Image {i+1})]: {ocr_result['text']}\n"
                                    
                                    # Store image data for display later
                                    extracted_image_data.append({
                                        'page': page_num + 1,
                                        'image_num': i + 1,
                                        'image': img,
                                        'ocr_text': ocr_result['text'],
                                        'ocr_details': ocr_result.get('details', [])
                                    })
                            except Exception as e:
                                st.warning(f"Error extracting text from image {i+1} on page {page_num+1}: {str(e)}")
            except Exception as e:
                st.warning(f"Warning: Could not extract text from page {page_num+1}: {str(e)}")
                continue
        
        # Store extracted image data in session state for later display
        if extracted_image_data:
            st.session_state['extracted_images'] = extracted_image_data
            st.success(f"✅ Found and processed {len(extracted_image_data)} images in the PDF document")
        else:
            st.session_state['extracted_images'] = []
            
        return text if text.strip() else "No text could be extracted from the PDF."
        
    except Exception as e:
        st.error(f"Error reading PDF: {str(e)}")
        return ""

def extract_text_from_image(image):
    """Extract text from an image using EasyOCR"""
    reader = load_ocr_reader()
    if reader is None:
        return "Error: OCR reader not available"
    
    try:
        # Convert PIL Image to numpy array
        if isinstance(image, Image.Image):
            image_np = np.array(image)
        else:
            image_np = image
            
        # Convert to RGB if image has alpha channel
        if len(image_np.shape) == 3 and image_np.shape[-1] == 4:
            image_np = cv2.cvtColor(image_np, cv2.COLOR_RGBA2RGB)
        
        # Get OCR results
        results = reader.readtext(image_np)
        
        # Extract text and confidence
        extracted_info = []
        for bbox, text, conf in results:
            extracted_info.append({
                'text': text,
                'confidence': conf,
                'bbox': bbox
            })
        
        return {
            'text': '\n'.join([info['text'] for info in extracted_info]),
            'details': extracted_info
        }
    except Exception as e:
        return f"Error processing image: {str(e)}"

def extract_images_from_pdf_page(file, page_num):
    """Extract images from a specific PDF page"""
    try:
        # Reset file pointer
        file.seek(0)
        
        # Load PDF
        pdf = PyPDF2.PdfReader(file)
        
        # Check if page exists
        if page_num >= len(pdf.pages):
            return []
            
        # Get the specified page
        page = pdf.pages[page_num]
        
        images = []
        if '/XObject' in page['/Resources']:
            xObject = page['/Resources']['/XObject'].get_object()
            
            for obj in xObject:
                if xObject[obj]['/Subtype'] == '/Image':
                    try:
                        data = xObject[obj].get_data()
                        img = Image.open(io.BytesIO(data))
                        
                        # Filter out very small images that are likely just icons or decorations
                        if img.width > 50 and img.height > 50:
                            images.append(img)
                    except:
                        continue
        
        return images
    except Exception as e:
        st.warning(f"Error extracting images from PDF page {page_num+1}: {e}")
        return []

def extract_images_from_pdf(file):
    """Extract all images from PDF file"""
    try:
        # Reset file pointer
        file.seek(0)
        
        # Load PDF
        pdf = PyPDF2.PdfReader(file)
        
        all_images = []
        for page_num in range(len(pdf.pages)):
            page_images = extract_images_from_pdf_page(file, page_num)
            for img in page_images:
                all_images.append({
                    'page': page_num + 1,
                    'image': img
                })
        
        return all_images
    except Exception as e:
        st.error(f"Error extracting images from PDF: {e}")
        return []

def extract_images_from_docx(file):
    """Extract images from Word document"""
    try:
        # Reset file pointer
        file.seek(0)
        
        # Create a temporary file to save the DOCX
        with open("temp.docx", "wb") as f:
            f.write(file.getvalue())
        
        # Get the relationship IDs for all images
        doc = docx.Document("temp.docx")
        image_rels = []
        for rel in doc.part.rels.values():
            if rel.reltype == RT.IMAGE:
                image_rels.append(rel)
                
        # Extract images using zipfile
        images = []
        with zipfile.ZipFile("temp.docx") as zip_ref:
            for rel in image_rels:
                # Get image URI
                image_part = rel.target_part
                image_data = image_part.blob
                
                try:
                    # Convert to PIL Image
                    img = Image.open(io.BytesIO(image_data))
                    
                    # Filter out very small images
                    if img.width > 50 and img.height > 50:
                        images.append(img)
                except Exception as e:
                    continue
                    
        # Clean up the temporary file
        os.remove("temp.docx")
        
        return images
    except Exception as e:
        st.error(f"Error extracting images from DOCX: {e}")
        if os.path.exists("temp.docx"):
            os.remove("temp.docx")
        return []

def extract_text_from_docx(file, extract_images_text=True):
    """Extract text from DOCX with optional image OCR"""
    try:
        # Create a bytes buffer
        file_bytes = io.BytesIO(file.getvalue())
        
        # Load the document
        doc = docx.Document(file_bytes)
        
        # Extract text from paragraphs
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
            
        # Process tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        # Extract and process images if requested
        if extract_images_text:
            # Extract images from DOCX
            images = extract_images_from_docx(file)
            
            # Process images with OCR
            extracted_image_data = []
            for i, img in enumerate(images):
                try:
                    ocr_result = extract_text_from_image(img)
                    if isinstance(ocr_result, dict) and ocr_result['text'].strip():
                        # Add image OCR text to the document text
                        text += f"\n[IMAGE TEXT (Image {i+1})]: {ocr_result['text']}\n"
                        
                        # Store image data for display
                        extracted_image_data.append({
                            'image_num': i + 1,
                            'image': img,
                            'ocr_text': ocr_result['text'],
                            'ocr_details': ocr_result.get('details', [])
                        })
                except Exception as e:
                    st.warning(f"Error extracting text from image {i+1} in DOCX: {str(e)}")
            
            # Store extracted image data in session state for later display
            if extracted_image_data:
                st.session_state['extracted_images'] = extracted_image_data
                st.success(f"✅ Found and processed {len(extracted_image_data)} images in the Word document")
            else:
                st.session_state['extracted_images'] = []
        
        return text
    except Exception as e:
        st.error(f"Error reading DOCX: {str(e)}")
        return ""

def extract_images_from_pptx(file):
    """Extract images from PowerPoint presentation"""
    try:
        # Reset file pointer
        file.seek(0)
        
        # Create a temporary file
        with open("temp.pptx", "wb") as f:
            f.write(file.getvalue())
        
        # Load presentation
        prs = pptx.Presentation("temp.pptx")
        
        images = []
        with zipfile.ZipFile("temp.pptx") as zip_ref:
            # Get all files in the zip
            for zip_info in zip_ref.infolist():
                # Check if it's an image file
                if zip_info.filename.startswith('ppt/media/') and any(zip_info.filename.endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif']):
                    try:
                        image_data = zip_ref.read(zip_info.filename)
                        img = Image.open(io.BytesIO(image_data))
                        
                        # Filter out very small images
                        if img.width > 50 and img.height > 50:
                            images.append(img)
                    except:
                        continue
                        
        # Clean up temporary file
        os.remove("temp.pptx")
        
        return images
    except Exception as e:
        st.error(f"Error extracting images from PPTX: {e}")
        if os.path.exists("temp.pptx"):
            os.remove("temp.pptx")
        return []

def extract_text_from_pptx(file, extract_images_text=True):
    """Extract text from PowerPoint with optional image OCR"""
    try:
        # Reset file pointer
        file.seek(0)
        
        # Load presentation
        prs = pptx.Presentation(file)
        
        # Extract text from slides
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            # Add slide number
            text += f"\n--- Slide {slide_num + 1} ---\n"
            
            # Get text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                    
        # Extract and process images if requested
        if extract_images_text:
            # Extract images from PPTX
            images = extract_images_from_pptx(file)
            
            # Process images with OCR
            extracted_image_data = []
            for i, img in enumerate(images):
                try:
                    ocr_result = extract_text_from_image(img)
                    if isinstance(ocr_result, dict) and ocr_result['text'].strip():
                        # Add image OCR text to the presentation text
                        text += f"\n[IMAGE TEXT (Image {i+1})]: {ocr_result['text']}\n"
                        
                        # Store image data for display
                        extracted_image_data.append({
                            'image_num': i + 1,
                            'image': img,
                            'ocr_text': ocr_result['text'],
                            'ocr_details': ocr_result.get('details', [])
                        })
                except Exception as e:
                    st.warning(f"Error extracting text from image {i+1} in PPTX: {str(e)}")
            
            # Store extracted image data in session state for later display
            if extracted_image_data:
                st.session_state['extracted_images'] = extracted_image_data
                st.success(f"✅ Found and processed {len(extracted_image_data)} images in the PowerPoint presentation")
            else:
                st.session_state['extracted_images'] = []
        
        return text
    except Exception as e:
        st.error(f"Error reading PPTX: {str(e)}")
        return ""

def extract_text_from_txt(file):
    return file.getvalue().decode("utf-8")

def extract_text_from_msg(file):
    # Save file temporarily
    with open("temp_email.msg", "wb") as f:
        f.write(file.getvalue())
    
    # Extract text using extract_msg
    msg = extract_msg.Message("temp_email.msg")
    text = f"Subject: {msg.subject}\n\nFrom: {msg.sender}\n\nTo: {msg.to}\n\nDate: {msg.date}\n\nBody:\n{msg.body}"
    
    # Remove temporary file
    os.remove("temp_email.msg")
    
    return text

def preprocess_text(text):
    # Tokenize and lowercase
    tokens = word_tokenize(text.lower())
    # Remove stopwords and punctuation
    stop_words = set(stopwords.words('english'))
    tokens = [token for token in tokens if token.isalnum() and token not in stop_words]
    # Lemmatize
    lemmatizer = WordNetLemmatizer()
    tokens = [lemmatizer.lemmatize(token) for token in tokens]
    return tokens

def generate_summary(text, sentences_count=5):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer()
    summary = summarizer(parser.document, sentences_count)
    return " ".join([str(sentence) for sentence in summary])

def extract_entities(text):
    if nlp is None:
        return {}
    
    doc = nlp(text)
    entities = {}
    
    for ent in doc.ents:
        if ent.label_ not in entities:
            entities[ent.label_] = []
        entities[ent.label_].append(ent.text)
    
    # Get unique entities
    for label in entities:
        entities[label] = list(set(entities[label]))
    
    return entities

def extract_key_phrases(text, n=10):
    tokens = preprocess_text(text)
    word_freq = Counter(tokens)
    # Get top n most common words
    return word_freq.most_common(n)

def analyze_sentiment(text, chunks=10):
    from textblob import TextBlob
    
    sentences = sent_tokenize(text)
    n = max(1, len(sentences) // chunks)
    chunks = [" ".join(sentences[i:i+n]) for i in range(0, len(sentences), n)]
    
    sentiments = []
    for chunk in chunks:
        blob = TextBlob(chunk)
        sentiments.append(blob.sentiment.polarity)
    
    return sentiments

def setup_gemini(api_key):
    """Configure the Gemini API with the provided key"""
    genai.configure(api_key=api_key)
    try:
        # Initialize and return Gemini 2.0 Flash model
        model = genai.GenerativeModel('gemini-2.0-flash')
        return model
    except Exception as e:
        st.error(f"Error initializing Gemini API: {e}")
        return None

def gemini_analyze_document(model, text, analysis_type="summarize", max_tokens=4000):
    """Use Gemini to analyze document text based on the requested analysis type"""
    
    # Add specific instruction for custom queries to be concise
    if analysis_type == "custom":
        prompt = f"Provide a very concise 1-2 line response to this question about the document. Be direct and specific:\n\n{text}"
    elif analysis_type == "analyze_image_text":
        prompt = f"This text was extracted from images using OCR. Please analyze it, correct any obvious OCR errors, and identify the key information:\n\n{text}"
    else:
        # Existing prompts for other analysis types
        prompts = {
            "summarize": f"Provide a concise summary of the following document in at most 5 paragraphs. Focus on the main ideas and key takeaways:\n\n{text}",
            "key_points": f"Extract and list the 10 most important key points from this document. Make each point clear and actionable:\n\n{text}",
            "action_items": f"Extract any action items, tasks, or to-dos mentioned in this document. Format them as a bulleted list with deadlines if available:\n\n{text}",
            "deep_analysis": f"Provide a deep analysis of this document covering its main themes, arguments, technical details, and conclusions. Include insights about the document's structure, tone, and intended audience:\n\n{text}",
            "extract_data": f"Extract structured information from this document including dates, numbers, statistics, financial figures, and relationships between entities. Format as a clearly organized list with categories:\n\n{text}"
        }
        prompt = prompts.get(analysis_type, prompts["summarize"])
    
    try:
        response = model.generate_content(prompt,
                                      generation_config=genai.types.GenerationConfig(
                                          max_output_tokens=max_tokens,
                                          temperature=0.2))
        return response.text
    except Exception as e:
        return f"Error processing with Gemini: {str(e)}"

def create_report_txt(file_name, text, analysis_results):
    """Create a formatted text report with all analysis results"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    report = f"""DOCUMENT ANALYSIS REPORT
Generated: {timestamp}
File: {file_name}

{'='*50}
SUMMARY
{'='*50}
{analysis_results.get('summary', 'No summary available')}

{'='*50}
KEY POINTS
{'='*50}
{analysis_results.get('key_points', 'No key points available')}

{'='*50}
ENTITIES EXTRACTED
{'='*50}
"""
    
    # Add entities
    entities = analysis_results.get('entities', {})
    if entities:
        for entity_type, items in entities.items():
            report += f"\n{entity_type}:\n"
            report += ", ".join(items) + "\n"
    else:
        report += "No entities extracted\n"
    
    # Add key phrases
    report += f"\n{'='*50}\nKEY PHRASES\n{'='*50}\n"
    key_phrases = analysis_results.get('key_phrases', [])
    if key_phrases:
        for phrase, count in key_phrases:
            report += f"{phrase}: {count}\n"
    else:
        report += "No key phrases extracted\n"
    
    # Add AI analysis if available
    if 'ai_analysis' in analysis_results:
        report += f"\n{'='*50}\nAI ANALYSIS\n{'='*50}\n"
        report += analysis_results['ai_analysis']
    
    # Add sentiment
    if 'sentiment' in analysis_results:
        report += f"\n{'='*50}\nSENTIMENT ANALYSIS\n{'='*50}\n"
        report += f"Overall sentiment: {analysis_results['sentiment']}\n"
    
    # Add statistics
    if 'statistics' in analysis_results:
        report += f"\n{'='*50}\nTEXT STATISTICS\n{'='*50}\n"
        stats = analysis_results['statistics']
        for key, value in stats.items():
            report += f"{key}: {value}\n"
            
    # Add OCR information if available
    if 'extracted_images' in analysis_results:
        report += f"\n{'='*50}\nEXTRACTED IMAGE TEXT\n{'='*50}\n"
        extracted_images = analysis_results['extracted_images']
        if extracted_images:
            for i, img_data in enumerate(extracted_images):
                report += f"\nImage {i+1}"
                if 'page' in img_data:
                    report += f" (Page {img_data['page']})"
                report += f":\n{img_data['ocr_text']}\n"
                report += f"Confidence: {sum([detail['confidence'] for detail in img_data['ocr_details']]) / len(img_data['ocr_details']):.2%}\n"
        else:
            report += "No image text extracted\n"
    
    return report

def display_extracted_images():
    """Display extracted images with OCR results in a gallery view"""
    if 'extracted_images' in st.session_state and st.session_state['extracted_images']:
        
        # Create a gallery of thumbnails
        st.subheader("📸 Extracted Images with OCR Text")
        
        # Use columns for the gallery
        num_cols = 3
        col_containers = st.columns(num_cols)
        
        # Organize images into the columns
        for i, img_data in enumerate(st.session_state['extracted_images']):
            col_idx = i % num_cols
            with col_containers[col_idx]:
                # Get the image
                img = img_data['image']
                
                # Convert PIL image to bytes for display
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                byte_img = buf.getvalue()
                
                # Create a unique key for the expander
                if 'page' in img_data:
                    expander_key = f"Page {img_data['page']}, Image {img_data['image_num']}"
                else:
                    expander_key = f"Image {img_data['image_num']}"
                
                # Create an expander for each image
                with st.expander(expander_key):
                    # Display the image
                    st.image(byte_img, caption=expander_key, use_column_width=True)
                    
                    # Display OCR text
                    st.markdown("**Extracted Text:**")
                    st.text_area("", img_data['ocr_text'], height=100, key=f"ocr_text_{i}")
                    
                    # Display confidence
                    if img_data['ocr_details']:
                        avg_confidence = sum([detail['confidence'] for detail in img_data['ocr_details']]) / len(img_data['ocr_details'])
                        confidence_color = "#10B981" if avg_confidence > 0.7 else "#F59E0B" if avg_confidence > 0.5 else "#EF4444"
                        
                        st.markdown(f"""
                        <div style="display: flex; align-items: center; margin-top: 10px;">
                            <div style="font-weight: 500; margin-right: 10px;">Confidence:</div>
                            <div style="background-color: {confidence_color}; color: white; padding: 3px 10px; border-radius: 15px;">
                                {avg_confidence:.2%}
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        # Display detailed OCR results
                        st.markdown("**Detailed OCR Results:**")
                        for j, detail in enumerate(img_data['ocr_details']):
                            st.markdown(f"""
                            <div style="padding: 5px; margin-bottom: 5px; border-bottom: 1px solid #eee;">
                                <div><strong>Text {j+1}:</strong> {detail['text']}</div>
                                <div><strong>Confidence:</strong> {detail['confidence']:.2%}</div>
                            </div>
                            """, unsafe_allow_html=True)
                    else:
                        st.info("No OCR details available for this image.")

def main():
    # Initialize session state for OCR settings
    if 'ocr_enabled' not in st.session_state:
        st.session_state['ocr_enabled'] = True
    
    # Sidebar for API key input and app info
    with st.sidebar:
        st.markdown("""
            <div style="display: flex; align-items: center; margin-bottom: 20px;">
                <div style="margin-right: 15px;">
                    <svg xmlns="http://www.w3.org/2000/svg" width="40" height="40" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
                        <path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"/>
                        <polyline points="14 2 14 8 20 8"/>
                        <circle cx="10" cy="13" r="2"/>
                        <path d="M18 19c-2.8 0-5-2.2-5-5"/>
                        <line x1="16" y1="13" x2="18" y2="13"/>
                    </svg>
                </div>
                <div>
                    <h1 style="margin: 0; color: #1E3A8A; font-size: 28px;">DocInsight Pro</h1>
                </div>
            </div>
        """, unsafe_allow_html=True)
        
        st.subheader("About the App")
        st.markdown("DocInsight Pro is an advanced document analyzer powered by Gemini 2.0 Flash AI.")

        st.markdown("---")
        st.subheader("Gemini AI Settings")
        api_key = st.text_input("Enter your Gemini API Key", type="password")
        
        if api_key:
            with st.spinner("Connecting to Gemini..."):
                st.session_state['gemini_model'] = setup_gemini(api_key)
                if st.session_state.get('gemini_model'):
                    st.success("✅ Gemini 2.0 Flash connected!")
                else:
                    st.error("❌ Connection failed. Check your API key.")
        else:
            st.info("ℹ️ Enter your Gemini API key to enable AI analysis.")
            st.session_state['gemini_model'] = None
        
        # OCR Settings
        st.markdown("---")
        st.subheader("OCR Settings")
        
        # Check if EasyOCR is available
        ocr_available = load_ocr_reader() is not None
        
        if ocr_available:
            st.success("✅ EasyOCR is enabled for image text extraction")
            
            # Toggle OCR for document images
            st.session_state['ocr_enabled'] = st.toggle("Extract text from document images", value=True)
            
            if st.session_state['ocr_enabled']:
                st.info("OCR will be performed on images within your documents.")
            else:
                st.info("OCR is disabled for document images.")
                
            # OCR language selection could be added here if needed
            # Currently using English as the default
        else:
            st.warning("⚠️ EasyOCR is not properly configured. Image text extraction may not work.")
            st.session_state['ocr_enabled'] = False
            
        st.markdown("---")
        st.caption("© 2025 DocInsight Pro | Version 1.0")
    
    # Main content - centralized header
    st.markdown('<h1 class="main-header">Document Analysis Dashboard</h1>', unsafe_allow_html=True)
    st.markdown('<p class="main-subheader">Upload, analyze, and extract key information from your documents using advanced AI.</p>', unsafe_allow_html=True)
    
    # File uploader with proper label
    uploaded_file = st.file_uploader(
        label="Upload Document",
        type=["pdf", "docx", "txt", "pptx", "msg", "png", "jpg", "jpeg"],
        key="document_uploader",
        help="Supported formats: PDF, Word, PowerPoint, Email, Text files, and Images"
    )
    
    if uploaded_file is not None:
        st.success(f"✅ File '{uploaded_file.name}' uploaded successfully")
        
        # Store file name for later use
        st.session_state['file_name'] = uploaded_file.name
        
        # Extract text based on file type
        file_type = uploaded_file.name.split(".")[-1].lower()
        
        with st.spinner("📝 Extracting text and processing images..."):
            if file_type in ["png", "jpg", "jpeg"]:
                # Handle image files
                image = Image.open(uploaded_file)
                
                # Display the image
                st.image(image, caption="Uploaded Image", use_column_width=True)
                
                # Extract text from image
                with st.spinner("Performing OCR..."):
                    ocr_result = extract_text_from_image(image)
                    if isinstance(ocr_result, dict):
                        text = ocr_result['text']
                        
                        # Store image data for display
                        st.session_state['extracted_images'] = [{
                            'image_num': 1,
                            'image': image,
                            'ocr_text': ocr_result['text'],
                            'ocr_details': ocr_result.get('details', [])
                        }]
                        
                        # Show detailed OCR results
                        with st.expander("View OCR Details"):
                            for idx, detail in enumerate(ocr_result['details']):
                                st.markdown(f"""
                                **Text {idx+1}:**
                                - Content: {detail['text']}
                                - Confidence: {detail['confidence']:.2%}
                                """)
                    else:
                        st.error(ocr_result)
                        text = ""
            elif file_type == "pdf":
                text = extract_text_from_pdf(uploaded_file, extract_images_text=st.session_state['ocr_enabled'])
            elif file_type == "docx":
                text = extract_text_from_docx(uploaded_file, extract_images_text=st.session_state['ocr_enabled'])
            elif file_type == "txt":
                text = extract_text_from_txt(uploaded_file)
            elif file_type == "pptx":
                text = extract_text_from_pptx(uploaded_file, extract_images_text=st.session_state['ocr_enabled'])
            elif file_type == "msg":
                text = extract_text_from_msg(uploaded_file)
            else:
                st.error("Unsupported file format")
                return
            
            # Store extracted text in session state
            st.session_state['document_text'] = text
            
            # Display extracted text
            with st.expander("View Extracted Text"):
                st.text_area("", text, height=200)
                
            # Display extracted images with OCR results, if any
            if 'extracted_images' in st.session_state and st.session_state['extracted_images']:
                display_extracted_images()
                
                # Store extracted images data for export
                st.session_state['analysis_results'] = st.session_state.get('analysis_results', {})
                st.session_state['analysis_results']['extracted_images'] = st.session_state['extracted_images']
        
        # Create container with shadow and rounded corners
        analysis_container = st.container()
        with analysis_container:
            # Document analysis options with improved tabs
            st.subheader("📊 Document Analysis")
            
            tabs = st.tabs([
                "📝 Basic Analysis", 
                "🧠 AI-Powered Analysis", 
                "🔍 Information Extraction", 
                "📤 Export Results"
            ])
            
            # To store all analysis results for export
            if 'analysis_results' not in st.session_state:
                st.session_state['analysis_results'] = {}
            
            with tabs[0]:
                # Document summarization
                st.subheader("Document Summary")
                
                # Skip summarization if text is too short
                if len(text.split()) < 10:
                    st.warning("The extracted text is too short for meaningful summarization.")
                    summary = "Text too short for summarization."
                    st.session_state['analysis_results']['summary'] = summary
                else:
                    summary_sentences = st.slider("Number of sentences", 3, 10, 5)
                    
                    with st.spinner("Generating summary..."):
                        summary = generate_summary(text, summary_sentences)
                        
                        # Add a nice card-like container for the summary
                        st.markdown(f"""
                        <div style="background-color: #F8F9FA; padding: 15px; border-radius: 5px; border-left: 5px solid #4F8BF9;">
                            {summary}
                        </div>
                        """, unsafe_allow_html=True)
                        
                        # Store summary for export
                        st.session_state['analysis_results']['summary'] = summary
                
                # Text statistics
                st.subheader("Text Statistics")
                words = word_tokenize(text)
                sentences = sent_tokenize(text)
                
                stats = {
                    "Total Words": len(words),
                    "Total Sentences": len(sentences),
                    "Avg Words per Sentence": round(len(words) / max(1, len(sentences)), 1),
                    "Estimated Reading Time": f"{round(len(words) / 200, 1)} min"  # Assuming 200 words per minute
                }
                
                # Create a nicer metric display
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Total Words", stats["Total Words"])
                with col2:
                    st.metric("Total Sentences", stats["Total Sentences"])
                with col3:
                    st.metric("Words per Sentence", stats["Avg Words per Sentence"])
                with col4:
                    st.metric("Reading Time", stats["Estimated Reading Time"])
                
                # Store statistics for export
                st.session_state['analysis_results']['statistics'] = stats
                
                # Sentiment Analysis
                st.subheader("Sentiment Analysis")
                with st.spinner("Analyzing sentiment..."):
                    sentiments = analyze_sentiment(text)
                    
                    # Create a nicer plot
                    fig, ax = plt.subplots(figsize=(10, 4))
                    plt.style.use('seaborn-v0_8')
                    ax.plot(sentiments, linewidth=2, color='#4F8BF9')
                    ax.fill_between(range(len(sentiments)), sentiments, alpha=0.2, color='#4F8BF9')
                    ax.set_title("Sentiment Throughout Document", fontsize=14, pad=20)
                    ax.set_ylabel("Sentiment (Negative to Positive)", fontsize=12)
                    ax.set_xlabel("Document Progress", fontsize=12)
                    ax.axhline(y=0, color='#EF4444', linestyle='--', alpha=0.5)
                    ax.spines['top'].set_visible(False)
                    ax.spines['right'].set_visible(False)
                    plt.tight_layout()
                    st.pyplot(fig)
                    
                    avg_sentiment = sum(sentiments) / len(sentiments)
                    sentiment_label = "Positive" if avg_sentiment > 0.05 else "Negative" if avg_sentiment < -0.05 else "Neutral"
                    sentiment_color = "#10B981" if avg_sentiment > 0.05 else "#EF4444" if avg_sentiment < -0.05 else "#F59E0B"
                    
                    # Display overall sentiment with colored badge
                    st.markdown(f"""
                    <div style="margin-top: 10px; margin-bottom: 20px;">
                        <span style="font-weight: 500;">Overall Document Sentiment:</span>
                        <span style="background-color: {sentiment_color}; color: white; padding: 3px 10px; border-radius: 15px; margin-left: 10px;">
                            {sentiment_label} ({avg_sentiment:.2f})
                        </span>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Store sentiment for export
                    st.session_state['analysis_results']['sentiment'] = f"{sentiment_label} ({avg_sentiment:.2f})"
            
            with tabs[1]:
                # Gemini-powered analysis
                st.subheader("AI-Powered Document Analysis")
                
                if not st.session_state.get('gemini_model'):
                    st.warning("⚠️ Enter your Gemini API key in the sidebar to enable AI analysis.")
                else:
                    # Add option to analyze image text if images were extracted
                    analysis_options = ["summarize", "key_points", "action_items", "deep_analysis", "extract_data"]
                    
                    # Add image text analysis option if images were extracted
                    if 'extracted_images' in st.session_state and st.session_state['extracted_images']:
                        analysis_options.append("analyze_image_text")
                    
                    analysis_type = st.selectbox(
                        "Select Analysis Type",
                        analysis_options,
                        format_func=lambda x: {
                            "summarize": "🔹 AI Summary",
                            "key_points": "🔹 Key Points",
                            "action_items": "🔹 Action Items & Tasks",
                            "deep_analysis": "🔹 Deep Analysis",
                            "extract_data": "🔹 Structured Data Extraction",
                            "analyze_image_text": "🔹 Image Text Analysis"
                        }[x]
                    )
                    
                    ai_col1, ai_col2 = st.columns([3, 1])
                    with ai_col2:
                        run_button = st.button("Run AI Analysis", type="primary")
                    
                    if run_button:
                        with st.spinner("🧠 Gemini 2.0 Flash is analyzing your document..."):
                            # If analyzing image text, gather all OCR text
                            if analysis_type == "analyze_image_text":
                                # Collect OCR text from all images
                                if 'extracted_images' in st.session_state and st.session_state['extracted_images']:
                                    ocr_text = ""
                                    for img_data in st.session_state['extracted_images']:
                                        if 'page' in img_data:
                                            ocr_text += f"[Image on Page {img_data['page']}]:\n{img_data['ocr_text']}\n\n"
                                        else:
                                            ocr_text += f"[Image {img_data['image_num']}]:\n{img_data['ocr_text']}\n\n"
                                    
                                    # Analyze OCR text
                                    result = gemini_analyze_document(st.session_state['gemini_model'], ocr_text, "analyze_image_text")
                                    st.session_state['analysis_results']['image_text_analysis'] = result
                                else:
                                    st.warning("No image text available for analysis.")
                                    result = "No image text available for analysis."
                            else:
                                result = gemini_analyze_document(st.session_state['gemini_model'], text, analysis_type)
                                
                                # Store AI analysis result for export
                                if analysis_type == "key_points":
                                    st.session_state['analysis_results']['key_points'] = result
                                else:
                                    st.session_state['analysis_results']['ai_analysis'] = result
                            
                            if analysis_type == "extract_data":
                                try:
                                    # Try to parse as JSON for better display
                                    # Remove markdown backticks if present
                                    json_text = re.sub(r'```json\s*|\s*```', '', result)
                                    json_data = json.loads(json_text)
                                    st.json(json_data)
                                except:
                                    st.markdown(result)
                            else:
                                st.markdown(result)
                    
                    # Customized AI analysis
                    st.subheader("Custom AI Query")
                    
                    # Add example queries dropdown
                    example_queries = [
                        "What are the main themes of this document?",
                        "Summarize the key arguments and evidence.",
                        "What are the potential weaknesses or gaps in this document?",
                        "Identify any technical terms and explain them.",
                        "What audience is this document intended for?",
                        "Extract all dates and events mentioned.",
                        "What are the next steps suggested in this document?"
                    ]
                    
                    # Add image-specific queries if images were extracted
                    if 'extracted_images' in st.session_state and st.session_state['extracted_images']:
                        example_queries.extend([
                            "What information is contained in the document images?",
                            "Summarize the text extracted from images.",
                            "Are there any discrepancies between the document text and image text?",
                            "What visual elements are described in the extracted image text?"
                        ])
                    
                    selected_example = st.selectbox("Choose an example query or write your own:", 
                                                   ["Custom"] + example_queries)
                    
                    if selected_example == "Custom":
                        custom_query = st.text_area("Enter your question about the document:", 
                                                  "What are the main themes of this document?",
                                                  height=80)
                    else:
                        custom_query = selected_example
                    
                    custom_col1, custom_col2 = st.columns([3, 1])
                    with custom_col2:
                        submit_button = st.button("Submit Query", type="primary")
                    
                    if submit_button:
                        with st.spinner("🔍 Processing your query..."):
                            prompt = f"Based on this document: \n\n{text}\n\n{custom_query}"
                            result = gemini_analyze_document(st.session_state['gemini_model'], 
                                                          prompt, 
                                                          "custom",
                                                          max_tokens=4000)
                            
                            # Create styled card for results
                            st.markdown("Query Result:")
                            
                            st.markdown(result)
                            
                            # Store custom query result
                            st.session_state['analysis_results']['custom_query'] = {
                                'query': custom_query,
                                'result': result
                            }
            
            with tabs[2]:
                # Information Extraction
                
                # Named Entity Recognition with improved visuals
                st.subheader("Named Entity Recognition")
                with st.spinner("Extracting entities..."):
                    entities = extract_entities(text)
                    
                    # Store entities for export
                    st.session_state['analysis_results']['entities'] = entities
                    
                    if entities:
                        # Create color mapping for entity types
                        entity_colors = {
                            "PERSON": "#3B82F6",      # Blue
                            "ORG": "#8B5CF6",         # Purple
                            "GPE": "#EC4899",         # Pink
                            "LOC": "#F97316",         # Orange
                            "DATE": "#14B8A6",        # Teal
                            "TIME": "#06B6D4",        # Cyan  
                            "MONEY": "#22C55E",       # Green
                            "PERCENT": "#EAB308",     # Yellow
                            "CARDINAL": "#64748B",    # Gray-blue
                            "ORDINAL": "#78716C",     # Gray
                            "QUANTITY": "#A855F7",    # Purple
                            "WORK_OF_ART": "#EC4899", # Pink
                            "LAW": "#6366F1",         # Indigo
                            "LANGUAGE": "#F97316",    # Orange
                        }
                        
                        # Get default color for other entity types
                        default_color = "#6B7280"  # Gray
                        
                        # Create cards for each entity type
                        cols = st.columns(3)
                        i = 0
                        
                        for entity_type, items in entities.items():
                            color = entity_colors.get(entity_type, default_color)
                            
                            with cols[i % 3]:
                                st.markdown(f"""
                                <div style="background-color: white; padding: 15px; border-radius: 8px; 
                                            margin-bottom: 15px; border-top: 4px solid {color}; box-shadow: 0 2px 4px rgba(0,0,0,0.05);">
                                    <p style="font-weight: 600; color: {color}; margin-bottom: 10px;">
                                        {entity_type} <span style="color: #6B7280; font-weight: 400;">({len(items)})</span>
                                    </p>
                                    <p style="font-size: 0.9rem;">{", ".join(items)}</p>
                                </div>
                                """, unsafe_allow_html=True)
                            
                            i += 1
                    else:
                        st.info("No entities found or spaCy model not loaded.")
                
                # Key phrases/topics with improved visualization
                st.subheader("Key Phrases")
                num_phrases = st.slider("Number of key phrases to extract", 5, 30, 10)
                
                with st.spinner("Extracting key phrases..."):
                    key_phrases = extract_key_phrases(text, num_phrases)
                    
                    # Store key phrases for export
                    st.session_state['analysis_results']['key_phrases'] = key_phrases
                    
                    # Create DataFrame for visualization
                    df = pd.DataFrame(key_phrases, columns=["Phrase", "Frequency"])
                    
                    # Create a nicer horizontal bar chart
                    plt.style.use('seaborn-v0_8')
                    fig, ax = plt.subplots(figsize=(10, 6))
                    
                    # Create custom colormap with vibrant colors
                    colors = [
                        '#FF6B6B',  # Coral Red
                        '#4ECDC4',  # Turquoise
                        '#45B7D1',  # Sky Blue
                        '#96CEB4',  # Sage Green
                        '#FFEEAD',  # Soft Yellow
                        '#D4A5A5',  # Dusty Rose
                        '#9B5DE5',  # Purple
                        '#4EA8DE',  # Ocean Blue
                        '#F15BB5',  # Pink
                        '#00BBF9',  # Bright Blue
                        '#00F5D4',  # Mint
                        '#FEE440',  # Yellow
                        '#8338EC',  # Royal Purple
                        '#FF006E',  # Hot Pink
                        '#FB5607'   # Orange
                    ]

                    # Repeat colors if there are more phrases than colors
                    colors = (colors * (len(df) // len(colors) + 1))[:len(df)]
                    
                    # Plot horizontal bars
                    bars = ax.barh(df["Phrase"], df["Frequency"], color=colors)
                    
                    # Add count labels at the end of each bar
                    for i, bar in enumerate(bars):
                        ax.text(bar.get_width() + 0.3, bar.get_y() + bar.get_height()/2, 
                                f"{df['Frequency'][i]}", va='center', fontsize=10)
                    
                    # Customize appearance
                    ax.set_title("Top Key Phrases", fontsize=14, pad=20)
                    ax.set_xlabel("Frequency", fontsize=12)
                    ax.invert_yaxis()  # To have highest count at the top
                    ax.spines['top'].set_visible(False)
                    ax.spines['right'].set_visible(False)
                    plt.tight_layout()
                    
                    st.pyplot(fig)
            
            with tabs[3]:
                # Export options
                st.subheader("Export Document Analysis Report")
                
                # Create TXT report
                txt_report = create_report_txt(
                    st.session_state.get('file_name', 'Unknown'),
                    st.session_state.get('document_text', ''),
                    st.session_state['analysis_results']
                )
                
                # Provide download option with styled button
                st.markdown("""
                <div style="text-align: center; margin: 30px 0;">
                """, unsafe_allow_html=True)
                
                # Single clear download button
                dl_button = st.download_button(
                    "📥 Download Complete Analysis Report",
                    txt_report,
                    file_name=f"docinsight_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain",
                    use_container_width=True
                )
                
                st.markdown("</div>", unsafe_allow_html=True)
                
                # Preview the TXT report
                with st.expander("📋 Preview Analysis Report"):
                    st.code(txt_report, language="text")
    else:
        # Instructions to get started
        with st.container():
            st.info("👆 Upload a document to get started with the analysis")
            
            # Main features in a cleaner layout
            st.markdown("## Key Features")
            
            feature_cols = st.columns(3)
            
            with feature_cols[0]:
                st.markdown("""
                <div class="feature-card">
                    <div class="icon-heading">
                        <svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path><polyline points="14 2 14 8 20 8"></polyline><line x1="16" y1="13" x2="8" y2="13"></line><line x1="16" y1="17" x2="8" y2="17"></line><polyline points="10 9 9 9 8 9"></polyline></svg>
                        <h3>Document Summarization</h3>
                    </div>
                    <p>Extract key insights and get concise summaries of your documents using both ML-based and AI-powered summarization techniques.</p>
                </div>
                """, unsafe_allow_html=True)
            
            with feature_cols[1]:
                st.markdown("""
                <div class="feature-card">
                    <div class="icon-heading">
                        <svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><circle cx="11" cy="11" r="8"></circle><line x1="21" y1="21" x2="16.65" y2="16.65"></line></svg>
                        <h3>Information Extraction</h3>
                    </div>
                    <p>Automatically detect entities, key phrases, and important information within your documents using advanced NLP algorithms.</p>
                </div>
                """, unsafe_allow_html=True)
            
            with feature_cols[2]:
                st.markdown("""
                <div class="feature-card">
                    <div class="icon-heading">
                        <svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 16V8a2 2 0 0 0-1-1.73l-7-4a2 2 0 0 0-2 0l-7 4A2 2 0 0 0 3 8v8a2 2 0 0 0 1 1.73l7 4a2 2 0 0 0 2 0l7-4A2 2 0 0 0 21 16z"></path><polyline points="3.27 6.96 12 12.01 20.73 6.96"></polyline><line x1="12" y1="22.08" x2="12" y2="12"></line></svg>
                        <h3>AI-Powered Deep Analysis</h3>
                    </div>
                    <p>Leverage Gemini 2.0 Flash to perform deep document analysis, extract structured data, and answer specific questions about your content.</p>
                </div>
                """, unsafe_allow_html=True)

        # Show supported document types in a clean grid
        st.markdown("## Supported Document Types")
        
        # Use 6 columns for the document types - added image column
        doc_cols = st.columns(6)
        
        # Document types with appropriate icons
        doc_types = [
            {
                "icon": '<svg xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path><polyline points="14 2 14 8 20 8"></polyline><line x1="16" y1="13" x2="8" y2="13"></line><line x1="16" y1="17" x2="8" y2="17"></line><polyline points="10 9 9 9 8 9"></polyline></svg>',
                "name": "Text Files (.txt)"
            },
            {
                "icon": '<svg xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path><polyline points="14 2 14 8 20 8"></polyline></svg>',
                "name": "PDF (.pdf)"
            },
            {
                "icon": '<svg xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path><polyline points="14 2 14 8 20 8"></polyline><rect x="8" y="12" width="8" height="6" rx="1"></rect></svg>',
                "name": "Word (.docx)"
            },
            {
                "icon": '<svg xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><rect x="2" y="3" width="20" height="14" rx="2" ry="2"></rect><line x1="8" y1="21" x2="16" y2="21"></line><line x1="12" y1="17" x2="12" y2="21"></line></svg>',
                "name": "PPT (.pptx)"
            },
            {
                "icon": '<svg xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M4 4h16c1.1 0 2 .9 2 2v12c0 1.1-.9 2-2 2H4c-1.1 0-2-.9-2-2V6c0-1.1.9-2 2-2z"></path><polyline points="22,6 12,13 2,6"></polyline></svg>',
                "name": "Email (.msg)"
            },
            {
                "icon": '<svg xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="#4F8BF9" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><rect x="3" y="3" width="18" height="18" rx="2" ry="2"></rect><circle cx="8.5" cy="8.5" r="1.5"></circle><polyline points="21 15 16 10 5 21"></polyline></svg>',
                "name": "Images (.png/.jpg)"
            }
        ]
        
        # Display document types in columns
        for i, doc_type in enumerate(doc_types):
            with doc_cols[i]:
                st.markdown(f"""
                <div class="document-type-card">
                    {doc_type["icon"]}
                    <h4 style="margin-top: 10px; color: #1E3A8A;">{doc_type["name"]}</h4>
                </div>
                """, unsafe_allow_html=True)
        
        # How it works section
        st.markdown("## How It Works")
        st.markdown("""
        <div style="background-color: white; padding: 20px; border-radius: 10px; box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);">
            <ol style="padding-left: 20px;">
                <li style="margin-bottom: 10px;"><strong>Upload</strong> your document (PDF, DOCX, PPTX, TXT, MSG, or Images)</li>
                <li style="margin-bottom: 10px;"><strong>Extract</strong> text and images with advanced OCR technology</li>
                <li style="margin-bottom: 10px;"><strong>Analyze</strong> your document using NLP and AI-powered techniques</li>
                <li style="margin-bottom: 10px;"><strong>Visualize</strong> insights including key phrases, entities, and sentiment</li>
                <li style="margin-bottom: 10px;"><strong>Export</strong> the complete analysis as a detailed report</li>
            </ol>
        </div>
        """, unsafe_allow_html=True)
        
        # Add a special highlight for new OCR feature
        st.markdown("""
        <div style="background-color: #EFF6FF; padding: 20px; border-radius: 10px; box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1); margin-top: 20px; border-left: 4px solid #3B82F6;">
            <h3 style="color: #1E3A8A; margin-top: 0;">✨ New Feature: Enhanced OCR</h3>
            <p style="margin-bottom: 0;">Our advanced OCR technology now extracts text from images embedded in your documents, 
            ensuring you don't miss any valuable information in charts, diagrams, or photos. Text from images is integrated
            seamlessly into your document analysis.</p>
        </div>
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()